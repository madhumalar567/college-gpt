import fitz
import time
import os
from flask import Flask, render_template, request, redirect, url_for, session, flash, jsonify
import sqlite3
from difflib import SequenceMatcher
from io import BytesIO
from PIL import Image
import pytesseract
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from google import genai
client = genai.Client(
    api_key="API_KEY"
)
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
current_study_text=""

app = Flask(__name__)
app.secret_key = "dsacoe_college_gpt_secure_session_key"

DATABASE_NAME = "college_gpt.db"

def get_db_connection():
    conn = sqlite3.connect(DATABASE_NAME)
    conn.row_factory = sqlite3.Row
    return conn


# ---------------- HOME PAGE ----------------
@app.route("/")
def home():

    if "reg_no" in session:
        return redirect(url_for("dashboard"))

    return render_template("index.html")


# ---------------- LOGIN ----------------
@app.route("/login", methods=["GET", "POST"])
def login():

    if "reg_no" in session:
        return redirect(url_for("dashboard"))

    if request.method == "POST":

        reg_no = request.form.get("registerNumber")
        password = request.form.get("password")

        if not reg_no or not password:
            flash("Please enter Register Number and Password.", "danger")
            return render_template("login.html")

        try:
            conn = get_db_connection()
            cursor = conn.cursor()

            cursor.execute(
                "SELECT * FROM students WHERE reg_no=? AND password=?",
                (reg_no, password)
            )

            student = cursor.fetchone()
            conn.close()

            if student:

                session["reg_no"] = student["reg_no"]
                session["name"] = student["name"]
                session["department"] = student["department"]
                session["year"] = student["year"]

                return redirect(url_for("dashboard"))

            else:
                flash("Invalid Register Number or Password.", "danger")

        except Exception:
            flash("Database Error!", "danger")

    return render_template("login.html")


# ---------------- DASHBOARD ----------------
@app.route("/dashboard")
def dashboard():

    if "reg_no" not in session:
        flash("Please login first.", "warning")
        return redirect(url_for("login"))

    return render_template("dashboard.html")

# ---------------- LOGOUT ----------------
@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("home"))

# ----------------------------------------------------
# CHAT MODULE ROUTES
# ----------------------------------------------------

@app.route("/chat")
def chat():
    """Renders the chat interface if the student is logged in."""
    if "reg_no" not in session:
        flash("Please log in to access the Chat Assistant.", "warning")
        return redirect(url_for("login"))
    return render_template("chat.html")


@app.route("/ask", methods=["POST"])
def ask():
    """Handles user queries and searches local SQLite tables based on keywords."""
    if "reg_no" not in session:
        return jsonify({"error": "Unauthorized"}), 401

    data = request.get_json()
    user_message = data.get("message", "").strip().lower()

    if not user_message:
        return jsonify({"response": "Please enter a valid question."})

    try:
        conn = get_db_connection()
        cursor = conn.cursor()

        # ---------------- HOD DETAILS ----------------
        if "hod" in user_message or "head of department" in user_message:
            department_map = {
                "it": "Information Technology",
                "information technology": "Information Technology",
                "cse": "Computer Science and Engineering",
                "computer science": "Computer Science and Engineering",
                "ece": "Electronics and Communication Engineering",
                "electronics": "Electronics and Communication Engineering",
                "eee": "Electrical and Electronics Engineering",
                "electrical": "Electrical and Electronics Engineering",
                "civil": "Civil Engineering",
                "mech": "Mechanical Engineering",
                "mechanical": "Mechanical Engineering",
                "ai": "Artificial Intelligence & Data Science",
                "aids": "Artificial Intelligence & Data Science",
                "artificial intelligence": "Artificial Intelligence & Data Science",
                "mba": "Management Studies (MBA)",
                "management": "Management Studies (MBA)",
                "physics": "Physics","chemistry": "Chemistry"
            }
            # If user asks only HOD Details, show all HODs
            if user_message.strip() in ["hod", "hod details", "head of department", "head of departments"]:
                cursor.execute("SELECT * FROM hods")
                rows = cursor.fetchall()
                response = "👨‍🏫 <b>All HOD Details</b><br><br>"
                for row in rows:
                    response += f"""
                    <b>{row['department']}</b><br>
                    👤 {row['hod_name']}<br>
                    📌 {row['designation']}<br>
                    📧 {row['email']}<br><br>"""
                conn.close()
                return jsonify({"response": response})
            selected_department = None
            for key, value in department_map.items():
                if key in user_message:
                    selected_department = value
                    break
            if selected_department is None:
                conn.close()
                return jsonify({
            "response": """
            👨‍🏫 <b>Please mention the department.</b><br><br>

            Examples:<br>
            • HOD of IT<br>
            • HOD of CSE<br>
            • HOD of ECE<br>
            • HOD of EEE<br>
            • HOD of Mechanical<br>
            • HOD of Civil<br>
            • HOD of AI & DS
            """
                })
            cursor.execute(
                "SELECT * FROM hods WHERE department=?",(selected_department,)
            )
            row = cursor.fetchone()
            conn.close()
            if row:
                return jsonify({
            "response": f"""
            👨‍🏫 <b>HOD Details</b><br><br>

            <b>Department :</b> {row['department']}<br>
            <b>Name :</b> {row['hod_name']}<br>
            <b>Designation :</b> {row['designation']}<br>
            <b>Specialization :</b> {row['specialization']}<br>
            <b>Email :</b> {row['email']}
            """
                })
            return jsonify({
        "response": "HOD details not found."
            })

        # -------- BUS ROUTES --------
        if "bus" in user_message or "route" in user_message:
            cursor.execute("SELECT * FROM bus_routes")
            bus_rows = cursor.fetchall()
            # ---------------- SHOW ALL ROUTES ----------------
            if user_message.strip() in [
                "bus",
                "bus details",
                "bus route",
                "bus routes",
                "college bus",
                "college buses"
            ]:
                response = "🚌 <b>Available College Bus Routes</b><br><br>"
                for row in bus_rows:
                    response += f"• <b>{row['route_no']}</b> - {row['route_name']}<br>"
                    response += """
                    <br><b>Examples:</b><br>
                    • Bus to Sathankulam<br>
                    • Tirunelveli Bus<br>
                    • Nazareth Route<br>
                    • Udangudi Bus<br>"""
                conn.close()
                return jsonify({"response": response})
            # ---------------- EXACT MATCH ----------------
            for row in bus_rows:
                route = row["route_name"].lower()
                if route in user_message:
                    conn.close()
                    return jsonify({
                "response": f"""
                🚌 <b>Bus Route Details</b><br><br>

                <b>Route No :</b> {row['route_no']}<br>
                <b>Route :</b> {row['route_name']}<br><br>

                <b>Stops :</b><br>
                {row['stops'].replace("->","<br>⬇<br>")}<br><br>

                <b>Timing :</b><br>
                {row['timing']}
                """
                    })
                for stop in row["stops"].lower().split("->"):
                    stop = stop.strip()
                    if stop in user_message:
                        conn.close()
                        return jsonify({
                    "response": f"""
                    🚌 <b>Bus Route Details</b><br><br>

                    <b>Route No :</b> {row['route_no']}<br>
                    <b>Route :</b> {row['route_name']}<br><br>

                    <b>Stops :</b><br>
                    {row['stops'].replace("->","<br>⬇<br>")}<br><br>

                    <b>Timing :</b><br>
                    {row['timing']}
                    """
                        })
            # ---------------- SIMILARITY MATCH ----------------
            best_match = None
            best_score = 0
            for row in bus_rows:
                score = SequenceMatcher(
                    None,
                    user_message,
                    row["route_name"].lower()
                    ).ratio()
                if score > best_score:
                    best_score = score
                    best_match = row
                for stop in row["stops"].lower().split("->"):
                    stop = stop.strip()
                    score = SequenceMatcher(
                        None,
                        user_message,
                        stop
                    ).ratio()
                    if score > best_score:
                        best_score = score
                        best_match = row
            if best_match and best_score > 0.45:
                conn.close()
                return jsonify({
            "response": f"""
            🚌 <b>Bus Route Details</b><br><br>

            <b>Route No :</b> {best_match['route_no']}<br>
            <b>Route :</b> {best_match['route_name']}<br><br>

            <b>Stops :</b><br>
            {best_match['stops'].replace("->","<br>⬇<br>")}<br><br>

            <b>Timing :</b><br>
            {best_match['timing']}
            """
                })
            conn.close()
            return jsonify({
        "response": """
        Sorry! Bus route not found.<br><br>

        Try like:<br>
        • Bus to Sathankulam<br>
        • Tirunelveli Bus<br>
        • Nazareth Route
        """
            })


        # ---------------- ADMISSION ----------------
        if "admission" in user_message or "fee" in user_message or "eligibility" in user_message or "apply" in user_message:

            cursor.execute("SELECT * FROM admissions")
            admission_rows = cursor.fetchall()

            for row in admission_rows:

                course = row["course"].lower()

                if course in user_message or "fee" in user_message or "admission" in user_message:

                    conn.close()

                    return jsonify({
                        "response": f"""
                        📋 <b>{row['course']}</b><br><br>
                        {row['details']}
                        """
                    })


        # ---------------- RULES ----------------
        if "rule" in user_message or "rules" in user_message or "attendance" in user_message or "discipline" in user_message:

            cursor.execute("SELECT * FROM rules")
            rule_rows = cursor.fetchall()

            response = "<b>College Rules</b><br><br>"

            for row in rule_rows:
                response += "• " + row["rule"] + "<br>"

            conn.close()

            return jsonify({
                "response": response
            })


        conn.close()

    except Exception as e:
        return jsonify({
            "response": f"Error : {str(e)}"
        })
    

    # -------- AI RESPONSE --------
    prompt=f"""
    You are CollegeGPT.
    You are the AI Assistant of Dr. Sivanthi Aditanar College of Engineering.
    Never say you are Gemini.
    Always introduce yourself as
    DSACOE CollegeGPT,
    AI Campus Assistant of
    Dr. Sivanthi Aditanar College of Engineering.
    Rules:
    1. Answer only educational questions.
    2. Maximum 120 words.
    3. Keep answers simple.
    4. Don't use markdown headings.
    5. Don't generate fake college information.
    6. If question is not related to DSACOE, answer normally.
    Student Question:{user_message}"""
    try:
        response = client.models.generate_content(
            model="gemini-flash-latest",
            contents=prompt
        )
        return jsonify({
            "response": response.text
        })
    except Exception as e:
        return jsonify({
            "response": str(e)
        })
# ----------------------------------------------------
# AI STUDY ASSISTANT MODULE ROUTES
# ----------------------------------------------------

@app.route("/study-assistant")
def study_assistant():
    """Renders the AI Study Workspace screen."""
    if "reg_no" not in session:
        flash("Please log in to access the Study Assistant workspace.", "warning")
        return redirect(url_for("login"))
    return render_template("study_assistant.html")

@app.route("/upload-notes", methods=["POST"])
def upload_notes():

    global current_study_text

    current_study_text = ""
    print("===== UPLOAD START =====")
    # ---------------- YOUTUBE ----------------
    youtube_url = request.form.get("youtube_url")

    if youtube_url:

        try:
            if "v=" in youtube_url:
                video_id = youtube_url.split("v=")[1].split("&")[0].split("?")[0]
            elif "youtu.be/" in youtube_url:
                video_id = youtube_url.split("youtu.be/")[1].split("?")[0]
            else:
                video_id = youtube_url.split("?")[0]
            print("Video ID:", video_id)
            api = YouTubeTranscriptApi()
            print("Fetching transcript...")
            transcript = api.fetch(video_id,languages=["ta", "en"])
            text = ""
            print("Transcript fetched successfully")
            for item in transcript:
                text += item.text + " "
            current_study_text = text
            return jsonify({
                "success": True,
                "message": "YouTube transcript uploaded successfully."
            })

        except Exception as e:
            print("YouTube Error:", e)
            return jsonify({
                "success": False,
                "message": str(e)
            })
        # ---------------- FILE ----------------

    if "file" not in request.files:

        return jsonify({
            "success": False,
            "message": "No file selected."
        })

    file = request.files["file"]

    if file.filename == "":

        return jsonify({
            "success": False,
            "message": "Please select a file."
        })
    filename = file.filename.lower()
    print("Filename :", filename)


    filepath = os.path.join(UPLOAD_FOLDER, file.filename)

    file.save(filepath)

    filename = file.filename.lower()
        # ---------------- PDF ----------------

    if filename.endswith(".pdf"):
        text = ""
        print("PDF detected")
        doc = fitz.open(filepath)
        for page in doc:
            page_text = page.get_text()
            print("Page Text Length :", len(page_text))
            print(repr(page_text[:100]))
            pix = page.get_pixmap(dpi=300)
            img = Image.open(BytesIO(pix.tobytes()))
            ocr_text = pytesseract.image_to_string(img)
            print("PyMuPDF Length :", len(page_text))
            print("OCR Length :", len(ocr_text))
            if len(page_text.strip()) > len(ocr_text.strip()):
                text += page_text + "\n"
            else:
                text += ocr_text + "\n"
        doc.close()
        current_study_text = text.strip()
        print("FINAL TEXT LENGTH =", len(current_study_text))
        print(current_study_text[:300])
        return jsonify({
            "success": True,
            "message": "PDF uploaded successfully."
        })
    
        # ---------------- PPT / PPTX ----------------

    if filename.endswith(".ppt") or filename.endswith(".pptx"):
        text = ""
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        current_study_text = text
        return jsonify({
            "success": True,
            "message": "PPT uploaded successfully."

        })
        # ---------------- IMAGE ----------------

    if filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
        print("IMAGE DETECTED")
        img = Image.open(filepath)
        text = pytesseract.image_to_string(img)
        print("OCR LENGTH:", len(text))
        current_study_text = text
        return jsonify({
            "success": True,
            "message": "Image uploaded successfully."
        })
    return jsonify({
        "success": False,
        "message": "Unsupported file format."
    })

@app.route("/generate-summary", methods=["POST"])
def generate_summary():
    global current_study_text
    print("Summary requested")
    print(len(current_study_text))
    if current_study_text.strip() == "":
        return jsonify({
            "summary": "Please upload a PDF, PPT, Image or YouTube lecture first."
        })
    prompt = f"""
You are an AI Study Assistant.
Use ONLY the uploaded study material below.
Do NOT use outside knowledge.
Create:
1. Easy Summary
2. Important Points
3. Key Concepts
4. Exam Tips
Study Material:
{current_study_text}
"""
    
    for i in range(3):
        try:
            response = client.models.generate_content(
                model="gemini-flash-latest",
                contents=prompt
            )

            return jsonify({
                "summary": response.text
            })

        except Exception as e:
            print(e)

            if "503" in str(e):
                time.sleep(5)
                continue

            return jsonify({
                "summary": str(e)
            })

    return jsonify({
        "summary": "Gemini server is busy. Please try again after a minute."
    })

@app.route("/generate-flashcards", methods=["POST"])
def generate_flashcards():

    global current_study_text

    if current_study_text.strip() == "":

        return jsonify({
            "raw_json": "Please upload study material first."
        })

    prompt = f"""
Generate EXACTLY 5 flashcards.

Return ONLY JSON.

Format:

[
 {{
   "front":"Question",
   "back":"Answer"
 }}
]

Study Material:

{current_study_text}
"""

    try:

        response = client.models.generate_content(

            model="gemini-flash-latest",

            contents=prompt

        )
        print(response.text)

        return jsonify({

            "raw_json": response.text

        })

    except Exception as e:
        print(e)
        return jsonify({
        "raw_json": str(e)
        })

@app.route("/generate-quiz", methods=["POST"])
def generate_quiz():

    global current_study_text

    if current_study_text.strip() == "":

        return jsonify({
            "raw_json":"Please upload study material first."
        })

    prompt = f"""
Generate 5 Multiple Choice Questions.

Return ONLY JSON.

Format:

[
 {{
   "id":1,
   "question":"Question",
   "options":[
      "Option A",
      "Option B",
      "Option C",
      "Option D"
   ],
   "answer":"Correct Option"
 }}
]

Use ONLY the uploaded notes.

Study Material:

{current_study_text}
"""

    try:

        response = client.models.generate_content(

            model="gemini-flash-latest",

            contents=prompt

        )
        print(response.text)
        return jsonify({

            "raw_json":response.text

        })

    except Exception as e:
        print(e)
        return jsonify({
        "raw_json": str(e)
        })
if __name__ == "__main__":
    app.run(debug=True)